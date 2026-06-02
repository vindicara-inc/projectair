"""Generate the Project AIR investor pitch deck as .pptx."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# Brand colors -- light airy palette (warm cream + soft lilac accents)
CREAM_BG = RGBColor(0xFA, 0xF7, 0xF2)         # warm cream slide background
LILAC = RGBColor(0x8B, 0x6C, 0xB0)            # soft lilac accent
LILAC_DARK = RGBColor(0x6B, 0x4F, 0x8E)       # deeper lilac for emphasis
SLATE_900 = RGBColor(0x1E, 0x1B, 0x2E)        # near-black for titles
SLATE_700 = RGBColor(0x3D, 0x38, 0x55)        # dark slate for body text
SLATE_500 = RGBColor(0x64, 0x5E, 0x78)        # medium for secondary text
SLATE_400 = RGBColor(0x89, 0x84, 0x9B)        # lighter secondary
WARM_WHITE = RGBColor(0xFF, 0xFF, 0xFF)        # white for stat boxes
STAT_BOX_BG = RGBColor(0xF0, 0xEC, 0xF5)      # very light lilac for stat boxes
STAT_BOX_BORDER = RGBColor(0xD8, 0xD0, 0xE4)  # lilac border
RED_ACCENT = RGBColor(0xC0, 0x3A, 0x4A)       # muted red for risk / alerts
TEAL = RGBColor(0x2A, 0x8D, 0x8D)             # teal for status labels

# Legacy aliases used throughout the slide code
BLACK = SLATE_900
DARK_BG = CREAM_BG
WHITE = SLATE_900
ZINC_300 = SLATE_700
ZINC_400 = SLATE_500
ZINC_500 = SLATE_400
RED = LILAC
CYAN = TEAL

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, font_name=BODY_FONT, alignment=PP_ALIGN.LEFT,
                 line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=ZINC_300, font_name=BODY_FONT, bullet_color=RED,
                    line_spacing=1.5):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.6)
        p.line_spacing = Pt(font_size * line_spacing)
        p.level = 0
    return txBox


def add_section_label(slide, text, top=0.6):
    add_text_box(slide, 0.8, top, 4, 0.4, text.upper(),
                 font_size=11, color=RED, bold=True, font_name=MONO_FONT)


def add_slide_title(slide, text, top=1.1):
    add_text_box(slide, 0.8, top, 11.5, 1.0, text,
                 font_size=36, color=WHITE, bold=True)


def add_red_bar(slide, top, width=12.0):
    shape = slide.shapes.add_shape(
        1, Inches(0.65), Inches(top), Inches(width), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()


def add_stat_box(slide, left, top, width, height, number, label, num_color=SLATE_900):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = STAT_BOX_BG
    shape.line.color.rgb = STAT_BOX_BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(32)
    p.font.color.rgb = num_color
    p.font.bold = True
    p.font.name = MONO_FONT

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = ZINC_500
    p2.font.name = BODY_FONT
    p2.space_before = Pt(8)


# =====================================================================
# SLIDE 1: Title / Cover
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_text_box(slide, 0.8, 1.5, 11.5, 1.2,
             "Project AIR",
             font_size=56, color=WHITE, bold=True)
add_text_box(slide, 0.8, 2.8, 11.5, 0.8,
             "Evidence-grade infrastructure for accountable AI agents.",
             font_size=28, color=ZINC_400)
add_red_bar(slide, 3.8, 3.0)
add_text_box(slide, 0.8, 4.2, 11.5, 0.5,
             "Seed Round  |  $2-3M",
             font_size=20, color=ZINC_300, font_name=MONO_FONT)
add_text_box(slide, 0.8, 5.0, 11.5, 0.5,
             "Kevin Minn  |  Founder & CEO  |  Kevin.Minn@vindicara.io",
             font_size=16, color=ZINC_500)
add_text_box(slide, 0.8, 5.6, 11.5, 0.5,
             "Vindicara, Inc.  |  Los Angeles, CA  |  vindicara.io",
             font_size=14, color=ZINC_500)
add_text_box(slide, 0.8, 6.6, 11.5, 0.3,
             "Confidential. Do not distribute.",
             font_size=11, color=ZINC_500, font_name=MONO_FONT)

# =====================================================================
# SLIDE 2: Description
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Description")
add_slide_title(slide, "What Project AIR is")

add_text_box(slide, 0.8, 2.3, 11.5, 1.5,
             "Project AIR is the forensic governance layer for autonomous AI agents. "
             "Every agent decision is written as a Signed Intent Capsule: a cryptographically "
             "signed envelope binding the declared goal, constraints, and context to each "
             "execution cycle. Each capsule carries a BLAKE3 content hash and an Ed25519 "
             "signature (with opt-in ML-DSA-65 post-quantum signing), chained to the previous "
             "step and anchored to public Sigstore Rekor transparency logs.",
             font_size=18, color=ZINC_300)

add_text_box(slide, 0.8, 4.2, 11.5, 0.6,
             "The result is evidence that survives subpoena, survives the vendor, "
             "and survives the auditor's first question.",
             font_size=20, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 5.2, 11.5, 2.0, [
    "MIT-licensed CLI + SDK on PyPI (projectair 0.7.1 live)",
    "14 detectors: 10/10 OWASP Agentic + 3 OWASP LLM + 1 AIR-native",
    "Four-layer architecture: detection, verification, causal reasoning, containment + cross-agent trust",
    "Six framework integrations: LangChain, OpenAI, Anthropic, LlamaIndex, Gemini, Google ADK",
    "Post-quantum ready: ML-DSA-65 (FIPS 204) experimental opt-in",
], font_size=16, color=ZINC_400)

# =====================================================================
# SLIDE 3: Problem
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Problem")
add_slide_title(slide, "AI agents are going off-script in production.\nNobody can prove what happened.")

add_bullet_list(slide, 0.8, 2.8, 5.5, 4.0, [
    "16,200 AI security incidents in 2025 (+49% YoY, Pillar Security)",
    "73% of production AI deployments have prompt injection vulnerabilities (OWASP / Lakera)",
    "Only 14% of organizations ship AI agents with full security approval (PwC 2025)",
    "EU AI Act enforcement begins August 2, 2026: audit trails and post-market monitoring required",
    "Non-compliance penalty: up to 7% of global annual revenue",
], font_size=16, color=ZINC_300)

add_bullet_list(slide, 7.0, 2.8, 5.5, 4.0, [
    "Prevention tools (Lakera, NeMo Guardrails, Bedrock Guardrails) try to stop bad things from happening",
    "None of them tell you what actually happened when an agent ran",
    "None produce evidence an auditor, regulator, or insurance carrier can use",
    "None bind a high-stakes action to the authenticated human who authorized it",
    "When something goes wrong, there is no chain of custody",
], font_size=16, color=ZINC_300)

# =====================================================================
# SLIDE 4: Solution
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Solution")
add_slide_title(slide, "Four layers of evidence-grade infrastructure")

layers = [
    ("Layer 0", "Detection", "10/10 OWASP Agentic + 3 LLM + 1 AIR-native. 14 detectors running over every chain.", "Shipped"),
    ("Layer 1", "External Trust Anchor", "RFC 3161 timestamps + Sigstore Rekor transparency log. Independently verifiable by anyone.", "Shipped (0.4.0)"),
    ("Layer 2", "Causal Reasoning", "Explains why each step happened. Hard edges go in the report; soft edges go in supporting context.", "Shipped (0.5.0)"),
    ("Layer 3", "Containment + Step-Up", "Halt agent actions. Require Auth0-verified human approval. Consent record, not just audit trail.", "Shipped (0.6.0)"),
]

for i, (layer, name, desc, status) in enumerate(layers):
    y = 2.5 + i * 1.2
    add_text_box(slide, 0.8, y, 1.5, 0.4, layer,
                 font_size=14, color=RED, bold=True, font_name=MONO_FONT)
    add_text_box(slide, 2.5, y, 3.5, 0.4, name,
                 font_size=18, color=WHITE, bold=True)
    add_text_box(slide, 2.5, y + 0.35, 6.5, 0.7, desc,
                 font_size=13, color=ZINC_400)
    add_text_box(slide, 10.5, y, 2.5, 0.4, status,
                 font_size=12, color=CYAN, font_name=MONO_FONT, alignment=PP_ALIGN.RIGHT)

# Layer 4 separate (alpha)
y = 2.5 + 4 * 1.2
add_text_box(slide, 0.8, y, 1.5, 0.4, "Layer 4",
             font_size=14, color=RED, bold=True, font_name=MONO_FONT)
add_text_box(slide, 2.5, y, 3.5, 0.4, "Cross-Agent Trust (A2A)",
             font_size=18, color=WHITE, bold=True)
add_text_box(slide, 2.5, y + 0.35, 6.5, 0.7,
             "AgDR Handoff Protocol. Cross-agent chain of custody with W3C Trace Context + Rekor counter-attestation.",
             font_size=13, color=ZINC_400)
add_text_box(slide, 10.5, y, 2.5, 0.4, "Shipped (0.7.0 alpha)",
             font_size=12, color=CYAN, font_name=MONO_FONT, alignment=PP_ALIGN.RIGHT)

# =====================================================================
# SLIDE 5: Team
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Team")
add_slide_title(slide, "Founder")

add_text_box(slide, 0.8, 2.5, 5.0, 0.5,
             "Kevin Minn", font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.8, 3.1, 5.0, 0.4,
             "Founder & CEO", font_size=18, color=CYAN)

add_bullet_list(slide, 0.8, 3.8, 5.5, 3.5, [
    "Built 6 AI products across 12+ months at SLTR Digital LLC",
    "Luminetic: AI-powered App Store compliance scanner (shipped)",
    "Cybersecurity domain expertise; full-time on Vindicara since March 2026",
    "Shipped 4-layer architecture, 14 detectors, 6 framework integrations, ops chain dogfood in < 90 days",
    "Solo founder velocity: 0.1.0 to 0.7.1 in 10 weeks",
], font_size=15, color=ZINC_300)

add_text_box(slide, 7.0, 2.5, 5.5, 0.5,
             "Planned hires (funded by raise)", font_size=20, color=WHITE, bold=True)
add_bullet_list(slide, 7.0, 3.3, 5.5, 3.5, [
    "Senior Engineer #1: Layer 1 Trust Anchor hardening, key rotation, key transparency",
    "Senior Engineer #2: Layer 4 handoff identity, Fulcio + OIDC cross-tenant federation",
    "Developer Relations: OSS adoption across the LangChain / agent framework ecosystem",
], font_size=15, color=ZINC_300)

# =====================================================================
# SLIDE 6: Market Size / Opportunity
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Market Size / Opportunity")
add_slide_title(slide, "The agent economy is here. Governance is not.")

add_stat_box(slide, 0.8, 2.5, 3.7, 1.8, "40%", "of enterprise apps will embed\ntask-specific AI agents by 2026\n(Gartner)")
add_stat_box(slide, 4.8, 2.5, 3.7, 1.8, "16,200", "AI security incidents in 2025\n+49% YoY\n(Pillar Security)", num_color=RED_ACCENT)
add_stat_box(slide, 8.8, 2.5, 3.7, 1.8, "Aug 2", "EU AI Act enforcement.\nAudit trails required.\nUp to 7% revenue penalty.", num_color=RED_ACCENT)

add_text_box(slide, 0.8, 4.8, 11.5, 0.6,
             "Every enterprise deploying AI agents will need forensic governance infrastructure. "
             "The question is whether they build it, buy it, or get caught without it.",
             font_size=18, color=ZINC_300)

add_bullet_list(slide, 0.8, 5.6, 5.5, 2.0, [
    "73% of production AI has prompt injection vulnerabilities (OWASP / Lakera)",
    "Only 14% ship agents with full security approval (PwC)",
    "Only 8% of MCP servers support OAuth (RSA 2026)",
], font_size=15, color=ZINC_400)

add_bullet_list(slide, 7.0, 5.6, 5.5, 2.0, [
    "MITRE ATLAS and NIST do not yet cover MCP attack vectors",
    "~50% of the agentic stack has zero standardized defensive guidance",
    "Bipartisan US political support for AI guardrails (Semafor, March 2026)",
], font_size=15, color=ZINC_400)

# =====================================================================
# SLIDE 7: Competitive Landscape
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Competitive Landscape")
add_slide_title(slide, "The last independent AI agent security platform.")

add_text_box(slide, 0.8, 2.3, 11.5, 0.5,
             "Acquirer consolidation created a vacuum. Vindicara fills it.",
             font_size=18, color=ZINC_400)

competitors = [
    ("Lakera", "Acquired by Check Point", "Guardrails only. No forensic chain. No causal reasoning. No containment."),
    ("CalypsoAI", "Acquired by F5 (gov-only)", "Government-locked. No developer self-serve. No open-source tier."),
    ("Guardrails AI", "$7.5M seed, 11 employees", "Input/output validation. No signed evidence. No cross-agent trust."),
    ("NeMo Guardrails", "NVIDIA OSS toolkit", "Complementary, not competitive. Partnership target (Tier 1 integration planned)."),
    ("Arize / Galileo", "Observability-first", "Tracing and evals, not forensic evidence. No compliance exports. No containment."),
]

for i, (name, status, diff) in enumerate(competitors):
    y = 3.0 + i * 0.85
    add_text_box(slide, 0.8, y, 2.5, 0.4, name,
                 font_size=16, color=WHITE, bold=True)
    add_text_box(slide, 3.5, y, 3.0, 0.4, status,
                 font_size=13, color=ZINC_500, font_name=MONO_FONT)
    add_text_box(slide, 7.0, y, 5.5, 0.4, diff,
                 font_size=13, color=ZINC_400)

add_red_bar(slide, 7.1, 11.7)
add_text_box(slide, 0.8, 7.2, 11.5, 0.5,
             "Vindicara: independent, developer-first, self-serve pricing, open-source top-of-funnel, full four-layer forensic stack.",
             font_size=14, color=WHITE, bold=True)

# =====================================================================
# SLIDE 8: Business Model
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Business Model")
add_slide_title(slide, "Open-core. MIT SDK top-of-funnel, commercial cloud + enterprise.")

tiers = [
    ("Open Source", "Free forever", "MIT CLI + SDK. 14 detectors. 6 framework integrations.\nSigned Intent Capsule chain. JSON/PDF/CEF exports.", "Live on PyPI"),
    ("Individual", "$39/mo ($350/yr)", "AIR Cloud client. Premium detectors + reports.\nNIST AI RMF, SOC2-AI templates. Email support.", "Stripe live"),
    ("Team", "$599/mo flat ($5,400/yr)", "Hosted AIR Cloud. Multi-agent dashboards. SIEM integrations\n(Datadog, Splunk, Sumo, Sentinel). Incident workflows.", "Stripe live"),
    ("Enterprise", "Custom ($50K+ ACV)", "SSO/SAML/RBAC. On-prem/VPC/air-gapped. Insurance carrier\nintegrations. SLA. Dedicated IR contact. BAA.", "Pipeline"),
]

for i, (tier, price, features, status) in enumerate(tiers):
    y = 2.5 + i * 1.2
    add_text_box(slide, 0.8, y, 2.2, 0.4, tier,
                 font_size=16, color=WHITE, bold=True)
    add_text_box(slide, 3.2, y, 2.5, 0.4, price,
                 font_size=16, color=CYAN, bold=True, font_name=MONO_FONT)
    add_text_box(slide, 6.0, y, 5.0, 0.8, features,
                 font_size=12, color=ZINC_400)
    add_text_box(slide, 11.5, y, 1.5, 0.4, status,
                 font_size=11, color=ZINC_500, font_name=MONO_FONT, alignment=PP_ALIGN.RIGHT)

# =====================================================================
# SLIDE 9: Traction
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Traction")
add_slide_title(slide, "Pre-revenue. Shipping velocity is the signal.")

add_stat_box(slide, 0.8, 2.5, 2.8, 1.5, "0.7.1", "Live on PyPI\n10 releases in 10 weeks")
add_stat_box(slide, 3.9, 2.5, 2.8, 1.5, "14", "Detectors shipping\n10/10 OWASP Agentic")
add_stat_box(slide, 7.0, 2.5, 2.8, 1.5, "4", "Layers shipped\nDetection to cross-agent trust")
add_stat_box(slide, 10.1, 2.5, 2.8, 1.5, "425+", "Tests passing\nmypy strict, ruff clean")

add_text_box(slide, 0.8, 4.5, 11.5, 0.5,
             "Verifiable proof, not claims:", font_size=20, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 5.1, 5.5, 2.5, [
    "Live Sigstore Rekor anchors: log index 1455601514 (Layer 1) and 1465403522 (Layer 4)",
    "Live Auth0 tenant integration demonstrated against production JWKS",
    "Vindicara dogfoods AIR on its own ops chain (vindicara.io/ops-chain)",
    "NVIDIA Inception Program member (April 2026)",
], font_size=15, color=ZINC_300)

add_bullet_list(slide, 7.0, 5.1, 5.5, 2.5, [
    "6 framework integrations: LangChain, OpenAI, Anthropic, LlamaIndex, Gemini, Google ADK",
    "Post-quantum signatures (ML-DSA-65 / FIPS 204) in working tree",
    "Full four-layer stack built by solo founder in < 90 days",
    "Stripe payment links live for Individual and Team tiers",
], font_size=15, color=ZINC_300)

# =====================================================================
# SLIDE 10: Go-to-Market Plan
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Go-to-Market Plan")
add_slide_title(slide, "Open-source wedge. Design partners. Enterprise contracts.")

phases = [
    ("Phase 1", "Now", "Public launch",
     "OSS on PyPI. HN + LinkedIn + founder outbound. 500 GitHub stars, 100 weekly pip installs in 90 days. 3-5 design partner conversations."),
    ("Phase 2", "Day 60", "First revenue",
     "AIR Cloud Team tier goes live. First design partner converts to paid. Target: first paying customer by day 60."),
    ("Phase 3", "Month 6", "Scale",
     "AIR Cloud GA. 10+ paying Team-tier teams. First enterprise LOI. SOC 2 Type I observation begins."),
    ("Phase 4", "Month 12-18", "Enterprise",
     "Enterprise tier live. 50 customers at $50K ACV. SOC 2 Type I complete. Series A pitch-ready."),
]

for i, (phase, timing, label, desc) in enumerate(phases):
    y = 2.5 + i * 1.2
    add_text_box(slide, 0.8, y, 1.5, 0.4, phase,
                 font_size=14, color=RED, bold=True, font_name=MONO_FONT)
    add_text_box(slide, 2.5, y, 1.5, 0.4, timing,
                 font_size=14, color=CYAN, font_name=MONO_FONT)
    add_text_box(slide, 4.2, y, 2.5, 0.4, label,
                 font_size=16, color=WHITE, bold=True)
    add_text_box(slide, 4.2, y + 0.35, 8.5, 0.7, desc,
                 font_size=13, color=ZINC_400)

add_text_box(slide, 0.8, 7.2, 11.5, 0.5,
             "Content-led: 'State of MCP Security' original research, EU AI Act Article 72 developer guide, quickstart tutorials.",
             font_size=13, color=ZINC_500)

# =====================================================================
# SLIDE 11: Product Roadmap
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Product Roadmap")
add_slide_title(slide, "Shipped fast. Building deeper.")

add_text_box(slide, 0.8, 2.3, 5.5, 0.5,
             "Shipped (last 10 weeks)", font_size=20, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 2.9, 5.5, 3.5, [
    "14 detectors, 10/10 OWASP Agentic coverage",
    "Layer 1: RFC 3161 + Sigstore Rekor anchoring",
    "Layer 2: Causal reasoning + forensic explanation",
    "Layer 3: Auth0-verified containment + step-up",
    "Layer 4 Wave 1: Cross-agent handoff protocol (alpha)",
    "ML-DSA-65 post-quantum signatures (experimental)",
    "6 framework integrations",
    "Ops chain dogfood on production infrastructure",
], font_size=14, color=ZINC_300)

add_text_box(slide, 7.0, 2.3, 5.5, 0.5,
             "Next 12 months", font_size=20, color=WHITE, bold=True)
add_bullet_list(slide, 7.0, 2.9, 5.5, 3.5, [
    "Layer 4 Wave 2: cross-tenant federation (Fulcio + OIDC)",
    "Layer 4 v1.5: Okta, Entra ID, SPIFFE enterprise adapters",
    "AIR Cloud: hosted ingestion + dashboards (Team tier)",
    "NVIDIA partnership Tier 1: NeMo Guardrails ingestion",
    "NVIDIA Tier 2: NemoGuard classification corroboration",
    "CrewAI, AutoGen, AG2 framework integrations",
    "Learned-baseline ASI10 anomaly detection (statistical profiling)",
    "Full ASI04 supply chain coverage beyond MCP",
], font_size=14, color=ZINC_300)

# =====================================================================
# SLIDE 12: Financial Summary
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Financial Summary")
add_slide_title(slide, "Path to $2.5M ARR in 18 months.")

add_text_box(slide, 0.8, 2.5, 11.5, 0.5,
             "Revenue model: land with OSS, convert to Team ($599/mo), expand to Enterprise ($50K+ ACV).",
             font_size=18, color=ZINC_300)

add_stat_box(slide, 0.8, 3.3, 3.7, 1.5, "50", "Target customers\nat $50K ACV in 18 months")
add_stat_box(slide, 4.8, 3.3, 3.7, 1.5, "$2.5M", "ARR target\n18-month horizon")
add_stat_box(slide, 8.8, 3.3, 3.7, 1.5, "$0", "Current revenue\nPre-revenue, pre-launch")

add_text_box(slide, 0.8, 5.3, 11.5, 0.5,
             "Unit economics (at scale):", font_size=18, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 5.9, 5.5, 2.0, [
    "OSS tier: $0 marginal cost (PyPI hosting is free)",
    "Team tier: ~95% gross margin (compute is minimal; chains are client-side)",
    "Enterprise: $50K-$250K ACV, contract-based",
], font_size=15, color=ZINC_300)

add_bullet_list(slide, 7.0, 5.9, 5.5, 2.0, [
    "AWS cost envelope: ~$5/mo at launch traffic (ops chain dogfood baseline)",
    "No GPU dependency; cryptographic operations run on CPU",
    "Revenue recognition: SaaS (monthly/annual subscriptions)",
], font_size=15, color=ZINC_300)

# =====================================================================
# SLIDE 13: Uses of Capital
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Uses of Capital")
add_slide_title(slide, "$2-3M seed. 18 months of runway.")

items = [
    ("Engineering (2 senior hires)", "60%",
     "Senior Engineer #1: Layer 1 Trust Anchor hardening, key rotation, key transparency log.\n"
     "Senior Engineer #2: Layer 4 handoff identity, Fulcio + OIDC cross-tenant federation."),
    ("Developer Relations", "15%",
     "Drive OSS adoption across the LangChain ecosystem. Content strategy: original research,\n"
     "EU AI Act guides, conference presence, community building."),
    ("Design Partner Conversion", "10%",
     "Convert 3-5 design partners to paid contracts. Fund pilot deployments, custom integrations,\n"
     "and security audits needed to close enterprise LOIs."),
    ("Infrastructure + Security", "10%",
     "AWS compute for AIR Cloud. SOC 2 Type I audit. Third-party security audit of the SDK.\n"
     "Domain, tooling, CI/CD."),
    ("Legal + Operations", "5%",
     "Delaware re-domicile. IP assignment. Standard SAFE/priced-round legal. Insurance."),
]

for i, (label, pct, desc) in enumerate(items):
    y = 2.4 + i * 1.0
    add_text_box(slide, 0.8, y, 5.0, 0.4, label,
                 font_size=16, color=WHITE, bold=True)
    add_text_box(slide, 6.0, y, 1.0, 0.4, pct,
                 font_size=16, color=CYAN, bold=True, font_name=MONO_FONT)
    add_text_box(slide, 7.2, y, 5.5, 0.8, desc,
                 font_size=12, color=ZINC_400)

# =====================================================================
# SLIDE 14: Risks
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Risks")
add_slide_title(slide, "Honest risk disclosure.")

risks = [
    ("Solo founder", "First hires contingent on raise. Mitigant: shipping velocity demonstrates solo-founder productivity; 4 layers + 14 detectors + 6 integrations in < 90 days."),
    ("Pre-revenue", "No paying customers yet. Mitigant: Stripe links live, pricing validated against market, regulatory forcing function (EU AI Act Aug 2) creates urgency."),
    ("Early-stage features", "Layer 4 is Wave 1 alpha. ML-DSA-65 is experimental. AIR Cloud dashboard is in development. Mitigant: core forensic chain and detection are production-grade."),
    ("Market timing", "Category is forming now. Large vendors (Cisco, NVIDIA, F5) are entering. Mitigant: Vindicara is the only independent, developer-first platform with self-serve pricing and full OWASP coverage."),
    ("Adoption risk", "OSS-to-paid conversion is uncertain. Mitigant: regulatory mandate creates must-have demand; open-core model proven by Snyk, GitLab, Elastic."),
]

for i, (risk, detail) in enumerate(risks):
    y = 2.4 + i * 0.95
    add_text_box(slide, 0.8, y, 2.5, 0.4, risk,
                 font_size=15, color=RED_ACCENT, bold=True)
    add_text_box(slide, 3.5, y, 9.0, 0.8, detail,
                 font_size=13, color=ZINC_400)

# =====================================================================
# SLIDE 15: Capital Efficiency
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "Capital Efficiency")
add_slide_title(slide, "Maximum output per dollar.")

add_text_box(slide, 0.8, 2.5, 11.5, 0.8,
             "Everything shipped to date was built by a solo founder in < 90 days, pre-raise, "
             "on personal runway. That includes a four-layer cryptographic architecture, "
             "14 production detectors, 6 framework integrations, a marketing site, live Stripe "
             "checkout, NVIDIA Inception membership, and a dogfooded ops chain on production infrastructure.",
             font_size=18, color=ZINC_300)

add_stat_box(slide, 0.8, 3.8, 3.7, 1.5, "$0", "External capital\nraised to date")
add_stat_box(slide, 4.8, 3.8, 3.7, 1.5, "< 90 days", "Time to ship\nfull four-layer stack")
add_stat_box(slide, 8.8, 3.8, 3.7, 1.5, "~$5/mo", "Current AWS cost\n(ops chain dogfood)")

add_text_box(slide, 0.8, 5.8, 11.5, 0.5,
             "What the raise unlocks:", font_size=20, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 6.3, 5.5, 1.5, [
    "Two senior engineers to harden L1 and ship L4 cross-tenant",
    "Developer relations for LangChain ecosystem OSS adoption",
    "Design partner conversion to 50 customers at $50K ACV",
], font_size=15, color=ZINC_300)

add_bullet_list(slide, 7.0, 6.3, 5.5, 1.5, [
    "SOC 2 Type I + third-party SDK security audit",
    "AIR Cloud hosted infrastructure (Team tier)",
    "18 months of runway at current burn trajectory",
], font_size=15, color=ZINC_300)

# =====================================================================
# SLIDE 16: Ask / Close
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.8, 2.0, 11.5, 1.2,
             "Project AIR",
             font_size=48, color=WHITE, bold=True)
add_text_box(slide, 0.8, 3.3, 11.5, 0.8,
             "Evidence-grade infrastructure for accountable AI agents.",
             font_size=24, color=ZINC_400)

add_red_bar(slide, 4.3, 3.0)

add_text_box(slide, 0.8, 4.7, 11.5, 0.5,
             "Raising $2-3M seed round", font_size=24, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 5.4, 11.5, 1.5, [
    "Two senior engineering hires (L1 hardening + L4 identity)",
    "Developer relations for OSS ecosystem adoption",
    "Design partner conversion: 50 customers at $50K ACV in 18 months",
], font_size=18, color=ZINC_300)

add_text_box(slide, 0.8, 6.5, 11.5, 0.5,
             "Kevin Minn  |  Kevin.Minn@vindicara.io  |  vindicara.io",
             font_size=16, color=ZINC_500, font_name=MONO_FONT)

# =====================================================================
# Save
# =====================================================================
out = Path("/Users/KMiI/Desktop/vindicara/docs/pitch/project-air-pitch-2026-05-10.pptx")
prs.save(str(out))
print(f"Saved to {out}")
